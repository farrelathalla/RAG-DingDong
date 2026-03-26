# backend/document_processor.py
from dataclasses import dataclass
from pathlib import Path
from typing import List
import fitz  # PyMuPDF
from backend.config import MAX_PAGE_CHARS

@dataclass
class PageEntry:
    doc_id: str
    doc_name: str
    page_num: int         # 1-based
    text: str

    def __post_init__(self):
        self.text = self.text[:MAX_PAGE_CHARS]


def extract_pages(file_path: str, doc_id: str, doc_name: str) -> List[PageEntry]:
    """Extract text per page from a document. Returns [] for unsupported formats."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(file_path, doc_id, doc_name)
    elif suffix == ".docx":
        return _extract_docx(file_path, doc_id, doc_name)
    elif suffix == ".pptx":
        return _extract_pptx(file_path, doc_id, doc_name)
    elif suffix in (".txt", ".md"):
        return _extract_text(file_path, doc_id, doc_name)
    else:
        return []


def _extract_pdf(file_path: str, doc_id: str, doc_name: str) -> List[PageEntry]:
    pages = []
    doc = fitz.open(file_path)
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            pages.append(PageEntry(doc_id=doc_id, doc_name=doc_name, page_num=i, text=text))
    doc.close()
    return pages


def _extract_docx(file_path: str, doc_id: str, doc_name: str) -> List[PageEntry]:
    from docx import Document
    doc = Document(file_path)
    PARAS_PER_PAGE = 30
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    pages = []
    for i in range(0, len(paragraphs), PARAS_PER_PAGE):
        chunk = "\n".join(paragraphs[i:i + PARAS_PER_PAGE])
        page_num = (i // PARAS_PER_PAGE) + 1
        pages.append(PageEntry(doc_id=doc_id, doc_name=doc_name, page_num=page_num, text=chunk))
    return pages


def _extract_pptx(file_path: str, doc_id: str, doc_name: str) -> List[PageEntry]:
    from pptx import Presentation
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if texts:
            pages.append(PageEntry(doc_id=doc_id, doc_name=doc_name, page_num=i, text="\n".join(texts)))
    return pages


def _extract_text(file_path: str, doc_id: str, doc_name: str) -> List[PageEntry]:
    text = Path(file_path).read_text(encoding="utf-8", errors="ignore").strip()
    if not text:
        return []
    return [PageEntry(doc_id=doc_id, doc_name=doc_name, page_num=1, text=text)]
